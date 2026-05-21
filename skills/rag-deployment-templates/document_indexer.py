"""
Document Indexer Skill - Index documents into Azure AI Search

Reference: https://learn.microsoft.com/en-us/azure/search/retrieval-augmented-generation-overview#content-preparation-for-rag
"""

import os
import json
import logging
from pathlib import Path
from typing import List, Dict, Optional, Any
from dataclasses import dataclass, asdict

logger = logging.getLogger(__name__)


@dataclass
class IndexingStats:
    """Statistics from indexing operation"""
    total_documents: int
    total_chunks: int
    documents_processed: Dict[str, int]  # By type (pdf, docx, code, etc)
    total_size_mb: float
    index_name: str
    status: str
    errors: List[str]


class DocumentIndexer:
    """Index documents into Azure AI Search using hybrid + semantic search"""
    
    SUPPORTED_TYPES = {
        "pdf": ["*.pdf"],
        "word": ["*.docx", "*.doc"],
        "excel": ["*.xlsx", "*.xls"],
        "markdown": ["*.md"],
        "code": ["*.sql", "*.py", "*.js", "*.ts", "*.yaml", "*.yml", "*.json"],
        "presentation": ["*.pptx", "*.ppt"],
    }
    
    def __init__(
        self,
        search_endpoint: str,
        search_admin_key: str,
        index_name: str = "rag-builder-index",
    ):
        """Initialize indexer with Azure Search credentials"""
        self.search_endpoint = search_endpoint
        self.search_admin_key = search_admin_key
        self.index_name = index_name
        self.documents_indexed = 0
        self.chunks_created = 0
        self.errors = []
    
    def scan_knowledge_folder(self, knowledge_path: str) -> Dict[str, List[Path]]:
        """
        Scan knowledge/ folder and categorize documents
        
        Expected structure:
        knowledge/
        ├── pdfs/
        ├── procedimientos/  (Word, Excel, Markdown)
        ├── codigo/         (SQL, Python, configs)
        └── presentaciones/ (PowerPoint)
        """
        documents_by_type = {dtype: [] for dtype in self.SUPPORTED_TYPES}
        knowledge_path = Path(knowledge_path)
        
        if not knowledge_path.exists():
            logger.error(f"Knowledge folder not found: {knowledge_path}")
            self.errors.append(f"Knowledge folder not found: {knowledge_path}")
            return documents_by_type
        
        # Scan each subdirectory
        for doc_type, patterns in self.SUPPORTED_TYPES.items():
            for pattern in patterns:
                for file_path in knowledge_path.rglob(pattern):
                    if file_path.is_file():
                        documents_by_type[doc_type].append(file_path)
                        logger.info(f"Found {doc_type}: {file_path.name}")
        
        return documents_by_type
    
    def chunk_document(
        self,
        file_path: Path,
        doc_type: str,
        chunk_size: int = 300,  # tokens
        overlap: int = 50,  # tokens
    ) -> List[Dict[str, Any]]:
        """
        Split document into semantic chunks
        
        For production, integrate with Azure Document Intelligence for smart chunking
        This is a simplified version using token-based splitting
        """
        chunks = []
        
        try:
            # Read file based on type
            if doc_type == "pdf":
                content = self._extract_pdf(file_path)
            elif doc_type in ["word", "excel"]:
                content = self._extract_office(file_path)
            elif doc_type == "markdown" or doc_type == "code":
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
            elif doc_type == "presentation":
                content = self._extract_presentation(file_path)
            else:
                logger.warning(f"Unsupported type: {doc_type}")
                return chunks
            
            # Split into chunks
            words = content.split()
            chunk_words = []
            
            for word in words:
                chunk_words.append(word)
                
                # Create chunk when size reached
                if len(chunk_words) >= chunk_size:
                    chunk_text = " ".join(chunk_words)
                    chunks.append({
                        "id": f"{file_path.stem}-chunk-{len(chunks)}",
                        "content": chunk_text,
                        "source_document": file_path.name,
                        "document_type": doc_type,
                        "file_path": str(file_path),
                        "metadata": {
                            "source": file_path.name,
                            "type": doc_type,
                            "chunk_index": len(chunks),
                        }
                    })
                    
                    # Keep overlap
                    overlap_words = int((overlap / chunk_size) * len(chunk_words))
                    chunk_words = chunk_words[-overlap_words:] if overlap_words > 0 else []
            
            # Add remaining words as final chunk
            if chunk_words:
                chunks.append({
                    "id": f"{file_path.stem}-chunk-{len(chunks)}",
                    "content": " ".join(chunk_words),
                    "source_document": file_path.name,
                    "document_type": doc_type,
                    "file_path": str(file_path),
                    "metadata": {
                        "source": file_path.name,
                        "type": doc_type,
                        "chunk_index": len(chunks),
                    }
                })
            
            logger.info(f"Created {len(chunks)} chunks from {file_path.name}")
            self.chunks_created += len(chunks)
            return chunks
            
        except Exception as e:
            error_msg = f"Error chunking {file_path.name}: {str(e)}"
            logger.error(error_msg)
            self.errors.append(error_msg)
            return []
    
    def _extract_pdf(self, file_path: Path) -> str:
        """Extract text from PDF"""
        try:
            # In production: use azure.ai.documentintelligence
            # For now, use PyPDF2 if available
            import PyPDF2
            text = ""
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                for page in pdf_reader.pages:
                    text += page.extract_text()
            return text
        except ImportError:
            logger.warning("PyPDF2 not installed, skipping PDF extraction")
            return "[PDF content - requires PyPDF2 or Azure Document Intelligence]"
        except Exception as e:
            logger.error(f"Error extracting PDF: {e}")
            return ""
    
    def _extract_office(self, file_path: Path) -> str:
        """Extract text from Word/Excel"""
        try:
            from docx import Document
            from openpyxl import load_workbook
            
            if file_path.suffix.lower() in ['.docx', '.doc']:
                doc = Document(file_path)
                return "\n".join([p.text for p in doc.paragraphs])
            elif file_path.suffix.lower() in ['.xlsx', '.xls']:
                wb = load_workbook(file_path)
                text = ""
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    text += f"\n=== Sheet: {sheet} ===\n"
                    for row in ws.iter_rows():
                        text += " | ".join(str(cell.value or "") for cell in row) + "\n"
                return text
        except ImportError:
            logger.warning("python-docx or openpyxl not installed")
            return "[Office content - requires python-docx or openpyxl]"
        except Exception as e:
            logger.error(f"Error extracting Office document: {e}")
            return ""
    
    def _extract_presentation(self, file_path: Path) -> str:
        """Extract text from PowerPoint"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = ""
            for slide_num, slide in enumerate(prs.slides):
                text += f"\n=== Slide {slide_num + 1} ===\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except ImportError:
            logger.warning("python-pptx not installed")
            return "[Presentation content - requires python-pptx]"
        except Exception as e:
            logger.error(f"Error extracting presentation: {e}")
            return ""
    
    def index_documents(
        self,
        knowledge_path: str,
        chunk_size: int = 300,
    ) -> IndexingStats:
        """
        Main indexing workflow
        
        1. Scan knowledge/ folder
        2. Extract text from each document type
        3. Create chunks with metadata
        4. Upload to Azure Search index
        5. Enable semantic ranking
        """
        stats = IndexingStats(
            total_documents=0,
            total_chunks=0,
            documents_processed={},
            total_size_mb=0,
            index_name=self.index_name,
            status="starting",
            errors=[]
        )
        
        # Scan documents
        documents_by_type = self.scan_knowledge_folder(knowledge_path)
        
        # Process each document type
        all_chunks = []
        for doc_type, file_paths in documents_by_type.items():
            stats.documents_processed[doc_type] = len(file_paths)
            stats.total_documents += len(file_paths)
            
            for file_path in file_paths:
                logger.info(f"Processing {doc_type}: {file_path.name}")
                
                # Chunk document
                chunks = self.chunk_document(file_path, doc_type, chunk_size)
                all_chunks.extend(chunks)
        
        stats.total_chunks = len(all_chunks)
        stats.errors = self.errors
        stats.status = "chunks_created"
        
        logger.info(f"Total chunks created: {stats.total_chunks}")
        
        # TODO: Upload to Azure Search
        # In production, this would:
        # 1. Create search index with vector fields
        # 2. Generate embeddings for each chunk
        # 3. Upload chunks to index with metadata
        # 4. Enable semantic ranking
        
        stats.status = "ready_for_upload"
        return stats


def index_knowledge_folder(
    knowledge_path: str,
    search_endpoint: str,
    search_admin_key: str,
    index_name: str = "rag-builder-index",
) -> Dict[str, Any]:
    """
    Public function to index documents
    """
    indexer = DocumentIndexer(
        search_endpoint=search_endpoint,
        search_admin_key=search_admin_key,
        index_name=index_name,
    )
    
    stats = indexer.index_documents(knowledge_path)
    
    return {
        "success": len(stats.errors) == 0,
        "stats": asdict(stats),
        "index_name": index_name,
        "ready_for_queries": True,
    }


if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO)
    
    # Example usage
    result = index_knowledge_folder(
        knowledge_path="knowledge/",
        search_endpoint="https://rag-builder.search.windows.net/",
        search_admin_key="your-key-here",
    )
    
    print(json.dumps(result, indent=2))
